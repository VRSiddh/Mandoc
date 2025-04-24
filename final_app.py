import streamlit as st
import os
import concurrent.futures
from io import BytesIO
from PIL import Image
from typing import Generator, Tuple, List, Optional
import numpy as np
import requests
from bs4 import BeautifulSoup
import fitz  # PyMuPDF
import docx
import pptx
import pytesseract
import faiss
import tempfile
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_core.embeddings import Embeddings
from langchain.embeddings import HuggingFaceBgeEmbeddings
from langchain.vectorstores import FAISS
from langchain_groq import ChatGroq
from langchain.prompts import PromptTemplate
from langchain.tools import tool
from langchain.chains import RetrievalQA
from langchain_community.document_loaders import CSVLoader, UnstructuredExcelLoader
from langchain_community.document_loaders.json_loader import JSONLoader
import pandas as pd


# Configuration
CHUNK_SIZE = 2000
CHUNK_OVERLAP = 200
MAX_WORKERS = os.cpu_count() or 8

# Set your API key here or use environment variable
os.environ["GROQ_API_KEY"] = "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"

# ====================
# Core Processing Pipeline
# ====================

class ParallelDocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            separators=["\n\n", "\n", ". "]
        )
        self.crawled_links = set()
        self._link_queue = []

    def process_files(self, file_objects: list) -> Generator[Document, None, None]:
        """Main entry point with parallel processing for Streamlit file objects"""
        # Create temporary files for processing
        temp_files = []
        for file in file_objects:
            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1])
            temp_file.write(file.getvalue())
            temp_file.close()
            temp_files.append((temp_file.name, file.name))  # Store original filename for better metadata
        
        # Process the temporary files
        with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            future_to_path = {
                executor.submit(self._process_single_file, fp, orig_name): fp 
                for fp, orig_name in temp_files
            }
            
            for future in concurrent.futures.as_completed(future_to_path):
                try:
                    yield from future.result()
                except Exception as e:
                    st.error(f"Error processing file: {e}")
        
        # Clean up temporary files
        for fp, _ in temp_files:
            try:
                os.unlink(fp)
            except:
                pass
    
    def _process_single_file(self, file_path: str, orig_filename: str) -> Generator[Document, None, None]:
        """Process individual files with format-specific handlers"""
        handler_map = {
            ".pdf": self._process_pdf,
            ".docx": self._process_docx,
            ".pptx": self._process_pptx,
            ".txt": self._process_txt,
            ".png": self._process_image,
            ".jpg": self._process_image,
            ".jpeg": self._process_image
        }
        
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == ".xlsx":
        # Try UnstructuredExcelLoader in 'elements' mode
            try:
                for doc in UnstructuredExcelLoader(file_path, mode='elements').load():
                    doc.metadata["source"] = orig_filename
                    yield doc
                return
            except Exception as e:
                st.warning(f"UnstructuredExcelLoader (elements mode) failed: {e}")
                # Try 'single' mode
                try:
                    for doc in UnstructuredExcelLoader(file_path, mode='single').load():
                        doc.metadata["source"] = orig_filename
                        yield doc
                    return
                except Exception as e2:
                    st.warning(f"UnstructuredExcelLoader (single mode) failed: {e2}")
                    # Fallback: Use pandas
                    try:
                        df = pd.read_excel(file_path)
                        text = df.to_string(index=False)
                        meta = {"source": orig_filename, "page": 1}
                        for chunk in self.text_splitter.split_text(text):
                            yield Document(page_content=chunk, metadata=meta)
                        return
                    except Exception as e3:
                        st.error(f"All Excel loaders failed: {e3}")
                        return
            
        if ext == ".json":
            try:
                for doc in JSONLoader(file_path, jq_schema='.').load():
                    doc.metadata["source"] = orig_filename  # Use original filename in metadata
                    yield doc
            except Exception as e:
                st.error(f"Error processing JSON file: {e}")
            return
            
        if ext == ".csv":
            try:
                for doc in CSVLoader(file_path, encoding='utf-8').load():
                    doc.metadata["source"] = orig_filename  # Use original filename in metadata
                    yield doc
            except Exception as e:
                st.error(f"Error processing CSV file: {e}")
            return
        
        handler = handler_map.get(ext)
        if not handler:
            st.warning(f"{ext} file handling not supported.")
            return
            
        flush_buffer = []
        flush_limit = CHUNK_SIZE * 2
        
        try:
            for page_content, metadata in handler(file_path):
                # Update metadata with original filename
                metadata["source"] = orig_filename
                
                flush_buffer.append((page_content, metadata))
                if sum(len(txt) for txt, _ in flush_buffer) >= flush_limit:
                    merged_text = "\n".join(txt for txt, _ in flush_buffer)
                    for chunk in self.text_splitter.split_text(merged_text):
                        yield Document(page_content=chunk, metadata=metadata.copy())
                    flush_buffer.clear()
            
            if flush_buffer:
                merged_text = "\n".join(txt for txt, _ in flush_buffer)
                for chunk in self.text_splitter.split_text(merged_text):
                    yield Document(page_content=chunk, metadata=flush_buffer[0][1].copy())
            
            if self._link_queue:
                with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
                    futures = [
                        executor.submit(self._crawl_and_process_link, link, parent_meta)
                        for link, parent_meta in self._link_queue
                        if link not in self.crawled_links
                    ]
                    for future in concurrent.futures.as_completed(futures):
                        try:
                            yield from future.result()
                        except Exception as e:
                            st.error(f"Error crawling link: {e}")
                self._link_queue.clear()
        except Exception as e:
            st.error(f"Error in document processing: {e}")
    
    # ====================
    # File Type Handlers (Optimized)
    # ====================
    
    def _process_pdf(self, file_path: str) -> Generator[Tuple[str, dict], None, None]:
        try:
            doc = fitz.open(file_path)
            
            with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
                text_futures = [
                    executor.submit(self._extract_pdf_page, doc, pn, file_path)
                    for pn in range(len(doc))
                ]
                
                for future in concurrent.futures.as_completed(text_futures):
                    page_num, text = future.result()
                    if text:
                        yield text, {"source": file_path, "page": page_num + 1}
            
            doc.close()
        except Exception as e:
            st.error(f"PDF processing error: {e}")
    
    def _extract_pdf_page(self, doc: fitz.Document, page_num: int, file_path: str) -> Tuple[int, str]:
        try:
            page = doc.load_page(page_num)
            text = page.get_text("text").strip()
            
            # If page has no text, try OCR
            if not text:
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text = pytesseract.image_to_string(img)
            
            # Extract links
            links = [link['uri'] for link in page.get_links() if link.get('kind') == fitz.LINK_URI and link.get('uri')]
            for link in links:
                self._link_queue.append((link, {"source": link, "parent": file_path, "page": page_num + 1}))
            
            # Extract and OCR images
            image_list = page.get_images(full=True)
            for img in image_list:
                try:
                    base_image = doc.extract_image(img[0])
                    image_bytes = base_image["image"]
                    img_text = self._image_to_text(image_bytes)
                    if img_text.strip():
                        text += "\n" + img_text
                except Exception as e:
                    st.error(f"Image OCR error: {e}")
            
            return page_num, text
        except Exception as e:
            st.error(f"Error reading page {page_num}: {e}")
            return page_num, ""
    
    def _process_docx(self, file_path: str) -> Generator[Tuple[str, dict], None, None]:
        try:
            doc = docx.Document(file_path)
            rels = doc.part.rels
            for para_num, para in enumerate(doc.paragraphs):
                content = []
                for run in para.runs:
                    el = run._element
                    for h in el.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}hyperlink'):
                        r_id = h.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                        if r_id and r_id in rels:
                            url = rels[r_id].target_ref
                            txt = ''.join(t.text for t in h.xpath('.//w:t', namespaces=el.nsmap) if t.text)
                            content.append(txt)
                            self._link_queue.append((url, {"source": url, "parent": file_path, "page": para_num + 1}))
                if para.text.strip():
                    content.append(para.text.strip())
                full_text = "\n".join(content).strip()
                if full_text:
                    yield full_text, {"source": file_path, "page": para_num + 1}
            
            for rel in doc.part._rels:
                rel_obj = doc.part._rels[rel]
                if "image" in rel_obj.target_ref:
                    image_bytes = rel_obj.target_part.blob
                    img_text = self._image_to_text(image_bytes)
                    if img_text.strip():
                        yield img_text, {"source": file_path, "page": "image"}
        except Exception as e:
            st.error(f"DOCX processing error: {e}")
    
    def _process_pptx(self, file_path: str) -> Generator[Tuple[str, dict], None, None]:
        try:
            prs = pptx.Presentation(file_path)
            texts = []
            links = []
            for slide_num, slide in enumerate(prs.slides):
                # Extract text from slide shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape.text.strip():
                        texts.append(shape.text.strip())
                    # Extract hyperlinks from shapes
                    if hasattr(shape, "hyperlink") and shape.hyperlink and shape.hyperlink.address:
                        links.append(shape.hyperlink.address)
                        self._link_queue.append(
                            (shape.hyperlink.address, {"source": shape.hyperlink.address, "parent": file_path, "page": slide_num + 1})
                        )
                # Extract text from notes
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    for shape in notes_slide.shapes:
                        if hasattr(shape, "text") and shape.text and shape.text.strip():
                            texts.append(shape.text.strip())
            # Yield all extracted text as a single chunk (will be split by splitter)
            if texts:
                yield "\n".join(texts), {"source": file_path, "page": "all"}
        except Exception as e:
            st.error(f"Error reading PPTX file: {e}")

    
    def _process_pptx_slide(self, slide, slide_num: int, file_path: str):
        content = []
        metadata = {"source": file_path, "page": slide_num + 1}
        
        def extract_links(shape):
            texts = []
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if hasattr(shape, "click_action") and shape.click_action and shape.click_action.hyperlink:
                url = shape.click_action.hyperlink.address
                if url:
                    self._link_queue.append((url, {"source": url, "parent": file_path, "page": slide_num + 1}))
            return texts
        
        for shape in slide.shapes:
            content.extend(extract_links(shape))
            if shape.shape_type == 13:  # Picture
                try:
                    image_bytes = shape.image.blob
                    img_text = self._image_to_text(image_bytes)
                    if img_text.strip():
                        content.append(img_text)
                except:
                    pass
        
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            for shape in notes_slide.shapes:
                content.extend(extract_links(shape))
        
        return [(("\n".join(content), metadata))] if content else []
    
    def _process_txt(self, file_path: str) -> Generator[Tuple[str, dict], None, None]:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
                if content.strip():
                    yield content, {"source": file_path, "page": 1}
        except Exception as e:
            st.error(f"TXT processing error: {e}")
    
    def _process_image(self, file_path: str) -> Generator[Tuple[str, dict], None, None]:
        try:
            with open(file_path, "rb") as f:
                image_bytes = f.read()
                text = self._image_to_text(image_bytes)
                if text.strip():
                    yield text, {"source": file_path, "page": 1}
        except Exception as e:
            st.error(f"Image processing error: {e}")
    
    def _image_to_text(self, image_bytes: bytes) -> str:
        try:
            image = Image.open(BytesIO(image_bytes))
            return pytesseract.image_to_string(image)
        except Exception as e:
            st.error(f"OCR error: {e}")
            return ""
    
    def _crawl_and_process_link(self, link: str, metadata: dict) -> Generator[Document, None, None]:
        if link in self.crawled_links:
            return
        self.crawled_links.add(link)
        
        try:
            response = requests.get(link, timeout=10)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, "html.parser")
            text = soup.get_text(separator="\n").strip()
            results = []
            for chunk in self.text_splitter.split_text(text):
                results.append(Document(
                    page_content=chunk,
                    metadata=metadata.copy()
                ))
            return results
        except Exception as e:
            st.warning(f"Failed to crawl {link}: {e}")
            return []


# ====================
# Optimized Vector Store Manager
# ====================

class VectorStoreManager:
    def __init__(self, 
                embeddings: Embeddings, 
                batch_size: int = 256,  # Tune based on GPU memory
                use_gpu: bool = False):  # Default to CPU for Streamlit
        self.embeddings = embeddings
        self.batch_size = batch_size
        self.index = None
        self._buffer = []
        
        if use_gpu:
            os.environ["FAISS_NO_AVX2"] = "1"  # Disable AVX for GPU compatibility
            self.res = faiss.StandardGpuResources()
    
    def add_documents(self, documents: List[Document]) -> None:
        """Batch-add documents with GPU acceleration"""
        self._buffer.extend(documents)
        
        if len(self._buffer) >= self.batch_size:
            self._flush_buffer()
    
    def finalize(self) -> FAISS:
        """Return finalized index"""
        if self._buffer:
            self._flush_buffer()
        return self.index
    
    def _flush_buffer(self) -> None:
        """Embed and add batched documents"""
        texts = [doc.page_content for doc in self._buffer]
        metadatas = [doc.metadata for doc in self._buffer]
        
        # Batch embed
        embeddings = self.embeddings.embed_documents(texts)
        
        # Convert to numpy array upfront
        embeddings_np = np.array(embeddings, dtype=np.float32)
        
        if self.index is None:
            # First batch: create index
            self.index = FAISS.from_embeddings(
                text_embeddings=list(zip(texts, embeddings)),
                embedding=self.embeddings,
                metadatas=metadatas
            )
            
            # Move to GPU if available
            if hasattr(self, 'res'):
                self.index.index = faiss.index_cpu_to_gpu(
                    self.res, 0, self.index.index
                )
        else:
            # Subsequent batches: add to existing index
            self.index.add_embeddings(
                text_embeddings=list(zip(texts, embeddings)),
                metadatas=metadatas
            )
        
        self._buffer.clear()


# ====================
# Retrieval and Q&A Functions
# ====================

@tool
def retrieve(query: str, vector_store, score_threshold: float = 0.7, k: int = 5) -> str:
    """
    Retrieves relevant data from the vector store
    
    Args:
        query: The search query
        vector_store: The FAISS vector store to search in
        score_threshold: Minimum similarity score to consider
        k: Number of results to return
    """
    try:
        results = vector_store.similarity_search_with_relevance_scores(query, score_threshold=score_threshold)
        if len(results) < k:
            results = vector_store.similarity_search_with_score(query, k=k)
        return results
    except:
        return []

def answer_question(query, vector_store):
    """Answer a question using the documents in the vector store"""
    try:
        # Initialize Groq LLM
        llm = ChatGroq(
        model="llama3-70b-8192",
        temperature=0.1,
        max_tokens=2048,
        )
        
        # Create QA chain with source tracking
        prompt_template = """Use the following context to answer the question.
        If you don't know the answer, say you don't know. Don't make up answers.
        Always cite your sources using the document metadata.

        Context:
        {context}

        Question: {question}

        Answer in markdown format with sources:
        """
        
        QA_PROMPT = PromptTemplate(
            template=prompt_template,
            input_variables=["context", "question"]
        )
        
        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            chain_type="stuff",
            retriever=vector_store.as_retriever(search_kwargs={"k": 5}),
            chain_type_kwargs={"prompt": QA_PROMPT},
            return_source_documents=True
        )
        
        result = qa_chain({"query": query})
        
        # Format response with sources
        response = {
            "answer": result["result"],
            "sources": list({f"{doc.metadata['source']} (page {doc.metadata['page']})"
                           for doc in result["source_documents"]})
        }
        return response
    except Exception as e:
        return {"answer": f"Error generating answer: {str(e)}", "sources": []}


# ====================
# Streamlit UI
# ====================

# Page config
st.set_page_config(page_title="Mandoc Chat", page_icon="🤖", layout="centered")

# Add custom CSS if needed
st.markdown("""
<style>
.chat-message {
    padding: 1.5rem; 
    border-radius: 0.5rem; 
    margin-bottom: 1rem;
    display: flex;
}
.chat-message.user {
    background-color: #2b313e;
}
.chat-message.assistant {
    background-color: #475063;
}
.chat-message .avatar {
    width: 20%;
}
.chat-message .content {
    width: 80%;
}
</style>
""", unsafe_allow_html=True)

# Sidebar for Uploads
st.sidebar.title("📁 Upload Documents")
uploaded_files = st.sidebar.file_uploader(
    "Supports PDF, DOCX, PPTX, XLSX, CSV, JSON, PNG, JPG, TXT",
    type=["pdf", "docx", "pptx", "xlsx", "csv", "json", "png", "jpg", "jpeg", "txt"],
    accept_multiple_files=True
)

# Initialize session state
if "chat" not in st.session_state:
    st.session_state.chat = []

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
    
if "processor" not in st.session_state:
    st.session_state.processor = ParallelDocumentProcessor()
    
if "embeddings" not in st.session_state:
    st.session_state.embeddings = HuggingFaceBgeEmbeddings(
        model_name="BAAI/bge-base-en",
        model_kwargs={'device': 'cpu'}  # Use CPU for Streamlit
    )
    
if "file_tracker" not in st.session_state:
    st.session_state.file_tracker = set()

# Main title
st.title("🤖 Mandoc Assistant")
st.caption("Upload documents and chat with your AI assistant.")

# Process uploaded files
if uploaded_files:
    # Check for new files
    current_files = {f.name for f in uploaded_files}
    new_files = [f for f in uploaded_files if f.name not in st.session_state.file_tracker]
    
    # Removed files
    removed_files = st.session_state.file_tracker - current_files
    if removed_files:
        st.session_state.vector_store = None  # Reset vector store if files were removed
        st.session_state.file_tracker = current_files
        st.info(f"Removed {len(removed_files)} files. Vector store has been reset.")
    
    # Process new files
    if new_files:
        with st.spinner(f"Processing {len(new_files)} new documents..."):
            # Initialize vector store if needed
            if st.session_state.vector_store is None:
                manager = VectorStoreManager(
                    embeddings=st.session_state.embeddings,
                    batch_size=256,
                    use_gpu=False  # Use CPU for Streamlit
                )
                
                # Process the documents
                doc_count = 0
                for doc in st.session_state.processor.process_files(new_files):
                    manager.add_documents([doc])
                    doc_count += 1
                
                # Finalize the vector store
                st.session_state.vector_store = manager.finalize()
                st.success(f"📄 Processed {doc_count} document chunks from {len(new_files)} files")
            else:
                # Add to existing vector store
                doc_count = 0
                for doc in st.session_state.processor.process_files(new_files):
                    st.session_state.vector_store.add_documents([doc])
                    doc_count += 1
                
                st.success(f"📄 Added {doc_count} document chunks from {len(new_files)} files")
            
            # Update file tracker
            for file in new_files:
                st.session_state.file_tracker.add(file.name)

# Display chat history
st.markdown("### Chat History")
for message in st.session_state.chat:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

# Chat input area
if prompt := st.chat_input("Ask a question about your documents..."):
    # Add user message to chat
    st.session_state.chat.append({"role": "user", "content": prompt})
    
    # Display the user message
    with st.chat_message("user"):
        st.markdown(prompt)
    
    # Generate a response if documents have been processed
    if st.session_state.vector_store:
        with st.spinner("Generating answer..."):
            # Get the answer from the Q&A function
            response = answer_question(prompt, st.session_state.vector_store)
            
            # Format the answer with sources
            if response["sources"]:
                answer_text = f"{response['answer']}\n\n**Sources:**\n"
                for source in response["sources"]:
                    answer_text += f"- {source}\n"
            else:
                answer_text = response["answer"]
                
            # Add to chat history and display
            st.session_state.chat.append({"role": "assistant", "content": answer_text})
            with st.chat_message("assistant"):
                st.markdown(answer_text)
    else:
        # No documents processed yet
        no_docs_message = "📚 Please upload documents first before asking questions."
        st.session_state.chat.append({"role": "assistant", "content": no_docs_message})
        with st.chat_message("assistant"):
            st.markdown(no_docs_message)

# Clear chat option
if st.session_state.chat:
    if st.sidebar.button("Clear Chat"):
        st.session_state.chat = []
        st.rerun()

# Reset everything option
if st.session_state.vector_store:
    if st.sidebar.button("Reset Documents & Chat"):
        st.session_state.vector_store = None
        st.session_state.chat = []
        st.session_state.file_tracker = set()
        st.rerun()

# Display some stats if documents are loaded
if st.session_state.vector_store:
    with st.sidebar.expander("Document Stats"):
        st.write(f"Loaded documents: {len(st.session_state.file_tracker)}")
        st.write(f"Document chunks: {len(st.session_state.vector_store.docstore._dict)}")

# Display some help information at the bottom
with st.expander("How to use Mandoc Assistant"):
    st.markdown("""
    1. **Upload Documents**: Use the sidebar to upload your documents (PDF, DOCX, PPTX, etc.)
    2. **Ask Questions**: Type your questions about the documents in the chat input
    3. **Get Answers**: The assistant will search through your documents and provide answers with sources
    4. **Manage Chat**: Use the sidebar buttons to clear the chat or reset everything
    """)
